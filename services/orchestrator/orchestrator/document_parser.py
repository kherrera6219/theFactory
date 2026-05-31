"""Document content extraction for multi-modal mission intake.

Extracts plain text from PDF / Word (.docx) / PowerPoint (.pptx) / Markdown /
plain-text attachments so the PM agent can reason over the *actual* document
content instead of just filename metadata.

All parsers are defensive: extraction failures (corrupt files, missing optional
dependencies, encrypted documents) return ``None`` and are logged at WARNING
level rather than raising — mission intake must proceed regardless.
"""
from __future__ import annotations

import io
import logging

LOGGER = logging.getLogger(__name__)

# Per-document extraction cap. Large documents are truncated so a single
# attachment cannot blow the LLM context window or the knowledge-lake payload.
MAX_EXTRACTED_CHARS = 50_000

_PDF_TYPES = ("application/pdf",)
_DOCX_TYPES = (
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
)
_PPTX_TYPES = (
    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
)
_TEXT_TYPES = ("text/markdown", "text/plain")


def parse_document(content: bytes, content_type: str, filename: str) -> str | None:
    """Extract text content from a document. Returns None if parsing fails.

    Routing is by ``content_type`` first, then by filename extension as a
    fallback (clients frequently send ``application/octet-stream``). Unknown
    formats return ``None`` so callers can fall back to metadata-only handling.
    """
    ctype = (content_type or "").split(";", 1)[0].strip().lower()
    name = (filename or "").lower()
    try:
        if ctype in _PDF_TYPES or name.endswith(".pdf"):
            return _parse_pdf(content)
        if ctype in _DOCX_TYPES or name.endswith(".docx"):
            return _parse_docx(content)
        if ctype in _PPTX_TYPES or name.endswith(".pptx"):
            return _parse_pptx(content)
        if ctype in _TEXT_TYPES or name.endswith((".md", ".txt")):
            return content.decode("utf-8", errors="replace")[:MAX_EXTRACTED_CHARS]
        return None
    except Exception as exc:  # noqa: BLE001 - intake must never fail on parse
        LOGGER.warning(
            "Document parse failed for %s: %s", filename, type(exc).__name__
        )
        return None


def _parse_pdf(content: bytes) -> str:
    import pypdf

    reader = pypdf.PdfReader(io.BytesIO(content))
    pages = [page.extract_text() or "" for page in reader.pages]
    return "\n\n".join(p for p in pages if p.strip())[:MAX_EXTRACTED_CHARS]


def _parse_docx(content: bytes) -> str:
    import docx

    doc = docx.Document(io.BytesIO(content))
    return "\n".join(
        para.text for para in doc.paragraphs if para.text.strip()
    )[:MAX_EXTRACTED_CHARS]


def _parse_pptx(content: bytes) -> str:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(content))
    texts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text)
    return "\n".join(texts)[:MAX_EXTRACTED_CHARS]
